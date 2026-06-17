#!/usr/bin/env python3
# 把 答辩PPT_png/ 里的 PNG 按文件名顺序拼成每页一张、铺满整页的 .pptx
# 用法: python3 build_pptx.py [PNG目录] [输出pptx]
# 例:  python3 build_pptx.py 答辩PPT_png 答辩PPT.pptx
import sys, glob, os
from pptx import Presentation
from pptx.util import Inches

PNG_DIR = sys.argv[1] if len(sys.argv) > 1 else "答辩PPT_png"
OUT      = sys.argv[2] if len(sys.argv) > 2 else "答辩PPT.pptx"

pngs = sorted(glob.glob(os.path.join(PNG_DIR, "*.png")))
if not pngs:
    sys.exit(f"在 {PNG_DIR}/ 没找到 PNG，请先运行 export_slides.mjs")

prs = Presentation()
# 16:9 宽屏（13.333 × 7.5 英寸），与 2560×1440 的 PNG 比例一致
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]  # 完全空白版式

for p in pngs:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(p, 0, 0, width=prs.slide_width, height=prs.slide_height)
    print(f"  ✓ {os.path.basename(p)}")

prs.save(OUT)
print(f"完成：{len(pngs)} 页 → {os.path.abspath(OUT)}")
